import traceback
from pathlib import Path
from functools import cached_property

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.utils.ppt_utils import duplicate_slide
from app.utils.ppt_common import create_text_chunks, find_and_replace_diagrams, print_shape_type, find_and_replace_OLE_photos, find_and_replace_OLE

class PresentationManager(object):
    """Contains Presentation object and functions to manage it"""
    
    # Character limit for content text in single slide
    MAX_CONTENT_LIMIT=2250

    def __init__(self, file_path=None, file=None, template_slide_index=1, slide_size=()):
        # Since presentation.Presentation class not intended to be constructed directly, using pptx.Presentation() to open presentation
        self.file_path = None
        
        if file:
            self.presentation = Presentation(file)
            print("Loaded presentation:", file.filename)
        elif file_path and Path(file_path).exists():
            self.presentation = Presentation(file_path)
            self.file_path = file_path
            print("Loaded presentation from:", file_path)
        else:
            self.presentation = Presentation()
            print("New presentation object loaded")

        if slide_size:
            height, width = slide_size
            self.set_slide_size(height, width)

        # Setting index of slide to be used as a template
        self.template_slide_index = template_slide_index

        # Get index of Blank slide layout
        layout_items_count = [len(layout.placeholders) for layout in self.presentation.slide_layouts]
        min_items = min(layout_items_count)
        self.blank_layout_id = layout_items_count.index(min_items)

        for slide in self.presentation.slides:
            find_and_replace_diagrams(slide) 
            find_and_replace_OLE_photos(slide)       
            # find_and_replace_OLE(slide)       

    @property
    def xml_slides(self):
        return self.presentation.slides._sldIdLst

    @property
    def _blank_slide_layout(self):        
        return self.presentation.slide_layouts[self.blank_layout_id]
    
    @property
    def total_slides(self):
        return len(self.presentation.slides)
    
    def set_slide_size(self, height, width):
        self.presentation.slide_height = height
        self.presentation.slide_width = width        

    def duplicate(self, index, destination=None):
        """
        Duplicates the slide with the given index. Adds slide to the end of the presentation
        """
        destination = destination or self
        try:
            slide = duplicate_slide(self.presentation, index, destination.presentation)
        except Exception:
            traceback.print_exc()

        return slide

    def move_slide(self, old_index, new_index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[old_index])
        self.xml_slides.insert(new_index, slides[old_index])

    def remove_slide(self, index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[index]) 

    def remove_all_slides(self):
        slides = list(self.xml_slides)
        for slide in slides:
            self.xml_slides.remove(slide)       
        

    def add_text_to_slide(self, index, text_content, title=""):
        """Adds title and content to slide at given index"""
        
        dest = self.presentation.slides[index]
        # Get title frame and content frame
        text_frames = []
        for shape in dest.shapes:
            if shape.has_text_frame:
                if 'Title' in shape.name:
                    title_frame = shape.text_frame
                else:
                    text_frames.append(shape.text_frame)
        # Choose first text frame as target
        content_frame = text_frames[0]
        
        # Clear content frame and add text
        content_frame.clear()
        p = content_frame.paragraphs[0]
        run = p.add_run()
        run.text = text_content

        # Clear title frame and add title
        title_frame.clear()
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = title

    def populate_slide(self, content, title=""):
        """Creates slides with given text and title, making more slides if text over limit"""
        
        duplicate_indices = []
        chunks = create_text_chunks(content, self.MAX_CONTENT_LIMIT)

        # Create slides for each chunk of text
        for chunk in chunks:
            slide_copy = self.duplicate(self.template_slide_index)
            i = self.presentation.slides.index(slide_copy)
            duplicate_indices.append(i)
            self.add_text_to_slide(i, chunk, title)

        # Move all slides to just before template slide
        new_index = self.template_slide_index + 1
        for old_index in duplicate_indices:
            self.move_slide(old_index, new_index)
            new_index += 1


    def save(self, filepath, remove_template=False):
        """Saves presentation to given filepath and removes slide used as template"""

        if remove_template:
            print("Removing template", self.template_slide_index)
            self.remove_slide(self.template_slide_index)
        self.presentation.save(filepath)
        print("Saved presentation to:", filepath)

    @classmethod
    def copy_slide_to_other_presentation(cls, source, dest_filepath, slides_to_copy=[]):
        # Load presentation
        destination = PresentationManager(dest_filepath)
        # Copy presentation size if destination is empty
        if destination.total_slides == 0:
            height, width = source.presentation.slide_height, source.presentation.slide_width
            destination.set_slide_size(height, width)
        # If no slide numbers given, default to all slides
        if not slides_to_copy:
            slides_to_copy = range(source.total_slides)

        for i in slides_to_copy:
            duplicate_slide(source.presentation, i, destination.presentation)
        # Save twice to avoid corruption bug
        destination.save(dest_filepath)
        destination = Presentation(dest_filepath)
        destination.save(dest_filepath)
   


    def _analyse_slide_elements(self, index, description=None):
        slide = self.presentation.slides[index]
        if description:
            print("*"*40, description, "*" * 40, sep="\n")
        for shape in slide.shapes:
            print_shape_type(shape)

    @cached_property
    def title(self):
        shapes = self.presentation.slides[0].shapes
        title = shapes.title.text if shapes.title else f"Untitled"
        return title

    def extract_all_text(self):
        docs = []
        slide_texts = []

        def get_text(shape):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = []
                for group_shape in shape.shapes:
                    text = get_text(group_shape)
                    if text:
                        group_text.append(text)
                return "\n".join(group_text)
            
            elif shape.has_text_frame:
                lines = shape.text.split("\n")
                text = "\n".join([line.strip() for line in lines if line.strip()])
                return text
            
            return ""

        for i, slide in enumerate(self.presentation.slides):
            shapes = slide.shapes
            title = shapes.title.text if shapes.title else f"Untitled Slide {i}"
            all_text = []
            for item in shapes:
                shape_text = get_text(item)
                if shape_text:
                    all_text.append(shape_text)
            slide_texts.append({
                "title": title,
                "content": "\n".join(all_text),
                "slide_id": slide.slide_id,
                "slide_index": i,
            })

        return slide_texts
    